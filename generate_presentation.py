"""
Run this script in your fatbox_env conda environment:
    pip install python-pptx
    python generate_presentation.py

Output: Fatbox_Afar_Presentation.pptx  (in the same folder as this script)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ───────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x37, 0x5E)   # slide backgrounds / titles
MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)   # accents, table headers
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)   # table alternating rows
ORANGE     = RGBColor(0xE8, 0x7A, 0x1E)   # highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout


# ── Helper functions ─────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=DARK_BLUE):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color=WHITE, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def accent_bar(slide, t=1.05, h=0.07, color=ORANGE):
    """Thin horizontal accent bar below the title."""
    box(slide, 0.4, t, 12.53, h, color=color)

def slide_title(slide, title, subtitle=None):
    """Standard title block used on content slides."""
    txt(slide, title, 0.4, 0.18, 12.0, 0.75,
        size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    accent_bar(slide)
    if subtitle:
        txt(slide, subtitle, 0.4, 1.18, 12.0, 0.4,
            size=16, color=RGBColor(0xCC, 0xDD, 0xEE), italic=True)

def bullet_block(slide, items, l, t, w, h,
                 size=17, color=WHITE, bullet="▸ ", line_gap=0.38):
    """Stack a list of strings as bullet points."""
    y = t
    for item in items:
        txt(slide, bullet + item, l, y, w, line_gap + 0.05,
            size=size, color=color)
        y += line_gap
    return y

def code_block(slide, code_lines, l, t, w):
    """Monospaced code box."""
    h = 0.32 * len(code_lines) + 0.15
    shape = box(slide, l, t, w, h, color=RGBColor(0x0D, 0x1B, 0x2A))
    txb = slide.shapes.add_textbox(Inches(l + 0.1), Inches(t + 0.05),
                                   Inches(w - 0.2), Inches(h - 0.1))
    txb.word_wrap = False
    tf = txb.text_frame
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xAA, 0xDD, 0xFF)
        run.font.name = "Courier New"
    return h

def table_2col(slide, rows, l, t, w,
               col_widths=(4.5, 7.5), row_h=0.38,
               header_color=MID_BLUE, alt_color=LIGHT_BLUE):
    """Simple 2-column table drawn as rectangles."""
    c1w, c2w = col_widths
    scale = w / (c1w + c2w)
    c1w, c2w = c1w * scale, c2w * scale

    y = t
    for i, (c1, c2) in enumerate(rows):
        is_header = (i == 0)
        bg_col = header_color if is_header else (LIGHT_BLUE if i % 2 == 0 else WHITE)
        txt_col = WHITE if is_header else DARK_GREY
        sz = 14 if is_header else 13
        bld = is_header

        box(slide, l,        y, c1w, row_h, color=bg_col)
        box(slide, l + c1w,  y, c2w, row_h, color=bg_col)

        txt(slide, c1, l + 0.08,       y + 0.04, c1w - 0.1, row_h - 0.05,
            size=sz, bold=bld, color=txt_col)
        txt(slide, c2, l + c1w + 0.08, y + 0.04, c2w - 0.1, row_h - 0.05,
            size=sz, bold=bld, color=txt_col)
        y += row_h
    return y - t


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
box(s, 0, 0, 13.33, 7.5, color=DARK_BLUE)
box(s, 0, 5.8, 13.33, 1.7, color=MID_BLUE)
box(s, 0, 3.5, 0.18, 2.2, color=ORANGE)

txt(s, "FATBOX", 0.5, 1.0, 12.0, 1.4,
    size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s, "Fault Analysis Toolbox", 0.5, 2.3, 12.0, 0.7,
    size=28, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)
txt(s, "Semi-automated fault extraction & structural analysis\nApplied to the Afar region — TanDEM-X 30 m DEM",
    0.5, 3.0, 12.0, 1.0,
    size=18, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT, italic=True)
txt(s, "Gayrin, Wrona & Brune (2025)  •  GFZ Helmholtz Centre for Geosciences",
    0.5, 6.05, 12.0, 0.5,
    size=14, color=WHITE, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is Fatbox?
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "What is Fatbox?",
            "An open-source Python toolbox for mapping and analysing fault networks")

# Left column — description
txt(s, "Fatbox extracts fault networks from various datasets and automatically "
       "computes their geometric and kinematic properties.",
    0.4, 1.55, 5.8, 1.0, size=16, color=WHITE)

bullets = [
    "~150 Python functions",
    "Works on DEMs, PIV strain data, numerical models",
    "Semi-automated extraction + full structural analysis",
    "Outputs: strike, throw, extension, displacement",
    "Jupyter notebook workflow — step by step",
]
bullet_block(s, bullets, 0.4, 2.55, 5.8, 3.0, size=15)

# Right column — 6 modules
box(s, 6.7, 1.35, 6.2, 5.7, color=RGBColor(0x12, 0x28, 0x4A))
txt(s, "The 6 modules", 6.9, 1.45, 5.8, 0.4,
    size=16, bold=True, color=ORANGE)

modules = [
    ("preprocessing.py",       "Load & prepare the DEM"),
    ("edits.py",               "Extract & edit the fault network"),
    ("metrics.py",             "Compute lengths and metrics"),
    ("plots.py",               "Visualise the network & results"),
    ("utils.py",               "Low-level helper functions"),
    ("structural_analysis.py", "Strike, dip, throw, extension…"),
]
y = 1.95
for name, desc in modules:
    txt(s, name, 6.9, y, 2.8, 0.35, size=13, bold=True,
        color=RGBColor(0xAA, 0xDD, 0xFF))
    txt(s, desc, 9.75, y, 2.9, 0.35, size=13, color=WHITE)
    y += 0.38


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Input data: the Afar DEM
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Input data — Afar DEM (TanDEM-X 30 m)")

# Info cards
cards = [
    ("File",        "TDM1_DEM_30m_UTM37.tif"),
    ("Projection",  "UTM Zone 37  (EPSG:32637)"),
    ("Resolution",  "30.2 m / pixel"),
    ("Size",        "7 343 × 11 073 pixels"),
    ("Coverage",    "~222 km (E–W)  ×  ~334 km (N–S)"),
    ("NW corner",   "X = 716 922 m   Y = 1 440 548 m  (UTM)"),
]
x_positions = [0.4, 4.55, 8.7]
for i, (label, val) in enumerate(cards):
    col = i % 3
    row = i // 3
    lx = x_positions[col]
    ly = 1.55 + row * 1.35
    box(s, lx, ly, 3.9, 1.2, color=MID_BLUE)
    box(s, lx, ly, 3.9, 0.38, color=RGBColor(0x1A, 0x50, 0x80))
    txt(s, label, lx + 0.12, ly + 0.05, 3.7, 0.32,
        size=13, bold=True, color=ORANGE)
    txt(s, val, lx + 0.12, ly + 0.42, 3.7, 0.65,
        size=15, color=WHITE)

# Coordinate conversion
box(s, 0.4, 4.3, 12.5, 1.55, color=RGBColor(0x0D, 0x1B, 0x2A))
txt(s, "Pixel ↔ UTM coordinate conversion", 0.6, 4.38, 12.0, 0.35,
    size=14, bold=True, color=ORANGE)
code_lines = [
    "UTM_Easting  = 716922 + col * 30.21          # column index → UTM X",
    "UTM_Northing = 1440548 - row * 30.21          # row index   → UTM Y",
    "col = (UTM_Easting  - 716922)  / 30.21        # UTM X → column",
    "row = (1440548 - UTM_Northing) / 30.21        # UTM Y → row",
]
y = 4.75
for line in code_lines:
    txt(s, line, 0.6, y, 12.2, 0.3, size=12,
        color=RGBColor(0xAA, 0xDD, 0xFF), italic=False)
    y += 0.27


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The workflow overview
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Workflow overview")

steps = [
    ("1", "Load DEM",             "Extract a tile with pixel\nor lat/lon coordinates"),
    ("2", "Fault extraction",     "Smooth → Canny edges\n→ Clean → Skeletonise"),
    ("3", "Classification",       "Build a graph: nodes =\npixels, edges = segments"),
    ("4", "Filtering",            "Connect, simplify, remove\nspurious components"),
    ("5", "Border processing",    "Remove edges whose cross-\nsection exits the DEM"),
    ("6", "Structural analysis",  "Strike, throw, extension,\ndisplacement per fault"),
    ("7", "Plots & export",       "Rose diagram, D/L scatter,\nextension map…"),
]

x = 0.25
for num, title, desc in steps:
    # Number circle
    box(s, x, 1.4, 0.55, 0.55, color=ORANGE)
    txt(s, num, x + 0.08, 1.42, 0.4, 0.45,
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Card
    box(s, x, 2.0, 1.55, 2.1, color=RGBColor(0x12, 0x28, 0x4A))
    txt(s, title, x + 0.08, 2.08, 1.4, 0.45,
        size=14, bold=True, color=ORANGE)
    txt(s, desc, x + 0.08, 2.58, 1.4, 1.5,
        size=12, color=WHITE)
    # Arrow (except last)
    if num != "7":
        txt(s, "→", x + 1.55, 1.55, 0.3, 0.4,
            size=24, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
    x += 1.84

txt(s, "All steps run in the Jupyter notebook  •  Each intermediate result can be saved and reloaded",
    0.4, 4.3, 12.5, 0.4, size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
    align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Fault extraction parameters
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Fault extraction — parameters",
            "Step 2: DEM  →  Smoothed  →  Edges  →  Cleaned  →  Skeleton")

# Pipeline arrow
pipeline = ["DEM", "Gaussian\nsmoothing", "Canny\nedges", "Remove\nnoise", "Skeleton"]
colors   = [MID_BLUE, RGBColor(0x2A,0x7A,0x5A), RGBColor(0x8A,0x3A,0x8A),
            RGBColor(0x9A,0x5A,0x10), ORANGE]
px = 0.3
for i, (label, col) in enumerate(zip(pipeline, colors)):
    box(s, px, 1.45, 2.1, 0.7, color=col)
    txt(s, label, px + 0.07, 1.5, 1.96, 0.6,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        txt(s, "→", px + 2.1, 1.6, 0.35, 0.4,
            size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    px += 2.45

# Parameter cards
params = [
    ("s_smoothed", "Gaussian σ — pre-filter",
     "Blurs the DEM before edge detection to suppress pixel noise.\n"
     "0 = no smoothing.  Typical: 1–3.\n"
     "↑ Increase if too many edges appear everywhere (lava, wadis)."),
    ("s_edges", "Canny σ — edge detection",
     "Controls the scale at which topographic discontinuities are detected.\n"
     "↓ Lower → fine features.  ↑ Higher → only strong broad scarps.\n"
     "Increase if rivers / lava flows are falsely detected."),
    ("msize_cleaned", "Minimum object size (pixels)",
     "Deletes connected groups smaller than this value.\n"
     "30 px × 30 m = 900 m minimum fault length.\n"
     "Afar faults are km-scale → can use 50–100 to remove noise."),
    ("con_cleaned", "Connectivity (leave at 2)",
     "Defines neighbour connectivity for noise removal.\n"
     "2 = 8-connectivity (includes diagonals).\n"
     "Do not change this parameter."),
]

x_pos = [0.3, 3.55, 6.8, 10.05]
for i, (name, subtitle, desc) in enumerate(params):
    lx = x_pos[i]
    box(s, lx, 2.35, 3.0, 4.7, color=RGBColor(0x12, 0x28, 0x4A))
    box(s, lx, 2.35, 3.0, 0.45, color=MID_BLUE)
    txt(s, name, lx + 0.1, 2.38, 2.8, 0.38,
        size=15, bold=True, color=WHITE)
    txt(s, subtitle, lx + 0.1, 2.85, 2.8, 0.35,
        size=12, bold=True, color=ORANGE)
    txt(s, desc, lx + 0.1, 3.28, 2.8, 3.6,
        size=12, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Tuning guide
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Tuning the extraction parameters",
            "Read the 6-panel overview plot and adjust until the skeleton matches real faults")

table_2col(s,
    [("Problem visible in the plot",          "Fix"),
     ("Too many tiny edges everywhere",        "↑  Increase  s_smoothed  and/or  msize_cleaned"),
     ("Real fault scarps not detected",        "↓  Decrease  s_edges  or  s_smoothed"),
     ("Rivers / lava flows detected as faults","↑  Increase  s_edges"),
     ("Short noise segments survive cleaning", "↑  Increase  msize_cleaned"),
     ("Long faults broken into many pieces",   "↓  Decrease  msize_cleaned  —  or use  connect_close_components  in filtering"),
    ],
    l=0.4, t=1.5, w=12.5,
    col_widths=(5.0, 7.5), row_h=0.52)

txt(s, "Recommended approach — hand mapping",
    0.4, 4.8, 12.5, 0.4, size=18, bold=True, color=ORANGE)

tips = [
    "Map a small representative tile in QGIS (10–20 km², a few hours of work)",
    "Export as a GeoTIFF raster at 30 m resolution (burn value = 1, nodata = 0)",
    "Load in Fatbox and overlay with the automated result",
    "Adjust parameters until the two networks match, then apply to the full DEM",
    "Pick a tile with both clear scarps AND noisy terrain (lava, wadis) for generalisability",
]
bullet_block(s, tips, 0.4, 5.25, 12.5, 2.5, size=14, line_gap=0.35)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Structural analysis
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Structural analysis — what is measured?",
            "A topographic cross-section is drawn perpendicular to each fault segment")

# Left: diagram description
box(s, 0.4, 1.5, 5.8, 5.5, color=RGBColor(0x12, 0x28, 0x4A))
txt(s, "Cross-section method", 0.55, 1.6, 5.5, 0.4,
    size=16, bold=True, color=ORANGE)
txt(s,
    "For each fault segment:\n"
    "  1.  Find the midpoint M of the segment\n"
    "  2.  Draw a profile perpendicular to the fault,\n"
    "       extending  d  pixels on each side of M\n"
    "  3.  Sample elevation along the profile\n"
    "  4.  Compute throw, dip, extension, displacement",
    0.55, 2.1, 5.5, 2.5, size=14, color=WHITE)

txt(s, "Key setting:  d = 12 pixels = 360 m each side",
    0.55, 4.65, 5.5, 0.4, size=14, bold=True, color=ORANGE)
txt(s,
    "use_natural_dip = False  →  dip_constrain = 60°\n"
    "(recommended when fault scarps are eroded)",
    0.55, 5.1, 5.5, 0.8, size=13, color=WHITE)

# Right: output quantities table
table_2col(s,
    [("Output",          "Meaning"),
     ("Strike",          "Azimuth direction of the fault (degrees)"),
     ("Throw",           "Max vertical offset across the scarp (metres)"),
     ("Natural dip",     "Apparent dip angle measured from the DEM"),
     ("Extension",       "Horizontal component of displacement (metres)"),
     ("Displacement",    "√(throw² + extension²)  —  total slip vector (m)"),
     ("Dip direction",   "Which side is the hanging wall (east / west)"),
    ],
    l=6.55, t=1.5, w=6.35,
    col_widths=(2.5, 3.85), row_h=0.48)

txt(s, "summary_matrix — one row per fault, 7 columns:",
    6.55, 5.0, 6.3, 0.35, size=13, bold=True, color=ORANGE)
code_lines = [
    "[ label | mean_strike° | length_m | max_dip° | mean_throw_m | mean_ext_m | max_disp_m ]",
]
code_block(s, code_lines, 6.55, 5.4, 6.35)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Output plots
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Output plots")

plots_list = [
    ("Rose diagram",         "Distribution of fault strike directions\n(azimuth histogram in polar form)"),
    ("Length histogram",     "Distribution of fault lengths\n(bars + cumulative dot plot)"),
    ("N ≥ length",           "Number of faults longer than x\n(power-law / exponential fit)"),
    ("Extension map",        "Edge attribute (extension / throw /\ndisplacement) coloured on the DEM"),
    ("D / L scatter",        "Max displacement vs length\nfor all faults in the network"),
    ("D / L profile",        "Displacement along a single fault\n(bell-shape = isolated fault growth)"),
]
x_pos = [0.3, 4.55, 8.8]
for i, (title, desc) in enumerate(plots_list):
    col = i % 3
    row = i // 3
    lx = x_pos[col]
    ly = 1.5 + row * 2.3
    box(s, lx, ly, 4.0, 2.1, color=RGBColor(0x12, 0x28, 0x4A))
    box(s, lx, ly, 4.0, 0.42, color=MID_BLUE)
    txt(s, title, lx + 0.12, ly + 0.06, 3.8, 0.34,
        size=15, bold=True, color=WHITE)
    txt(s, desc, lx + 0.12, ly + 0.55, 3.8, 1.4,
        size=13, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — The adapted Afar notebook
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Adapted notebook — Afar region",
            "tutorials/digital_elevation_models/Tuto_DEM_Afar_notebook.ipynb")

# Changes table
table_2col(s,
    [("What",               "Original tutorial → Adapted for Afar"),
     ("Google Colab cells", "Drive mount + pip install → Removed"),
     ("modules path",       "/content/drive/…/modules → /home/guiguizz/stage/modules"),
     ("save_path",          "Google Drive path → /home/guiguizz/stage/tutorials/…"),
     ("DEM loading",        "Pre-saved .npy arrays → Loads from TDM1_DEM_30m_UTM37.tif"),
     ("Coordinate system",  "'geographic' (lat/lon) → 'grid coordinates' (pixel indices)"),
     ("Default tile",       "Full 12 km² demo DEM → First 1 000 × 1 000 px (≈ 30 km²)"),
     ("Output filenames",   "*_GLO* → *_afar*"),
    ],
    l=0.4, t=1.5, w=12.5,
    col_widths=(3.8, 8.7), row_h=0.46)

txt(s, "How to pick your tile (pixel indices):",
    0.4, 5.55, 12.5, 0.35, size=15, bold=True, color=ORANGE)
code_block(s,
    ["coord_north = 0     # first row  (row 0 = northern edge of DEM)",
     "coord_south = 1000  # last row   (1 000 px × 30 m = 30 km south)",
     "coord_west  = 0     # first col  (col 0 = western edge of DEM)",
     "coord_east  = 1000  # last col   → for full DEM use 7343 / 11073"],
    l=0.4, t=5.95, w=12.5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — How to use hand mapping
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Using hand mapping to tune parameters",
            "Tuto_DEM_GLO30_hand_map_notebook.ipynb")

# Left: what the notebook does
box(s, 0.4, 1.5, 5.8, 5.5, color=RGBColor(0x12, 0x28, 0x4A))
txt(s, "What the hand-map notebook does", 0.55, 1.6, 5.5, 0.4,
    size=15, bold=True, color=ORANGE)
txt(s,
    "Instead of automated edge detection,\n"
    "load a GeoTIFF of manually-drawn faults\n"
    "(black/white raster: fault = 1, background = 0).\n\n"
    "The rest of the pipeline is identical:\n"
    "  • filter → border processing\n"
    "  • structural analysis\n"
    "  • same plots\n\n"
    "Also overlays both networks (hand vs auto)\n"
    "on the same hillshade to compare visually.",
    0.55, 2.1, 5.5, 4.7, size=14, color=WHITE)

# Right: QGIS workflow
box(s, 6.55, 1.5, 6.35, 5.5, color=RGBColor(0x12, 0x28, 0x4A))
txt(s, "How to create the hand-map GeoTIFF in QGIS",
    6.7, 1.6, 6.0, 0.4, size=15, bold=True, color=ORANGE)

qgis_steps = [
    ("1", "Draw fault traces as a multiline vector layer (shapefile)"),
    ("2", "Rasterize with GDAL 'Rasterize (vector to raster)':\n"
          "     Burn value = 1  •  Nodata = 0\n"
          "     Resolution = 30 m  •  Same extent as DEM tile"),
    ("3", "Export as GeoTIFF (same CRS as DEM)"),
    ("4", "Load in Fatbox with extract_geotiff_and_coord\n"
          "     using the same pixel coordinates as the DEM"),
    ("5", "Overlay and adjust s_smoothed, s_edges, msize_cleaned\n"
          "     until automated result matches hand map"),
]
y = 2.1
for num, step in qgis_steps:
    box(s, 6.7, y, 0.38, 0.38, color=ORANGE)
    txt(s, num, 6.72, y + 0.02, 0.35, 0.34,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, step, 7.15, y, 5.6, 0.75, size=12, color=WHITE)
    y += 0.82 if "\n" in step else 0.48


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Summary / Next steps
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK_BLUE)
slide_title(s, "Summary & next steps")

# Three columns
cols = [
    ("What we have", ORANGE, [
        "TanDEM-X 30 m DEM of the Afar region",
        "Fatbox installed and configured",
        "Adapted notebook ready to run",
        "Full workflow documented in FATBOX_NOTES.md",
    ]),
    ("Immediate next steps", MID_BLUE, [
        "Open Tuto_DEM_Afar_notebook.ipynb",
        "Choose a representative 30 km² tile",
        "Run the extraction and inspect the 6-panel plot",
        "Tune s_smoothed, s_edges, msize_cleaned",
        "Save the filtered network (save_G = True)",
    ]),
    ("For best results", RGBColor(0x2A,0x7A,0x5A), [
        "Map a small area by hand in QGIS",
        "Use hand-map notebook to compare networks",
        "Validate parameters before scaling to full DEM",
        "Use save/load at each stage to avoid re-running",
    ]),
]
x_pos = [0.3, 4.55, 8.8]
for i, (title, col, items) in enumerate(cols):
    lx = x_pos[i]
    box(s, lx, 1.5, 4.1, 0.5, color=col)
    txt(s, title, lx + 0.12, 1.56, 3.9, 0.38,
        size=16, bold=True, color=WHITE)
    box(s, lx, 2.0, 4.1, 5.0, color=RGBColor(0x12, 0x28, 0x4A))
    y = 2.1
    for item in items:
        txt(s, "▸  " + item, lx + 0.12, y, 3.85, 0.55, size=13, color=WHITE)
        y += 0.62

txt(s, "Gayrin, Wrona & Brune (2025)  •  https://doi.org/10.5281/zenodo.15716080",
    0.4, 7.1, 12.5, 0.3, size=12,
    color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER, italic=True)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "Fatbox_Afar_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
